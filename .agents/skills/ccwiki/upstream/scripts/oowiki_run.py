"""oowiki - LLM Wiki 지식 체계 유틸리티 스크립트

유틸리티 기능(status, lint, index, list)을 담당.
실제 수집·통합(run)은 Claude(AI)가 주도.
"""
import sys
import os
from pathlib import Path
from datetime import datetime

if sys.stdout and hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

WIKI_DIR = Path(__file__).parents[3].parent / "01_obsidian" / "0020_wiki"
INBOX_DIR = Path(__file__).parents[3].parent / "01_obsidian" / "0019_정리"

BINARY_EXTENSIONS = {
    # 이미지 (텍스트 추출 불가)
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".ico", ".tiff", ".tif", ".svg",
    # 영상/오디오 (텍스트 추출 불가)
    ".mp4", ".avi", ".mov", ".mkv", ".mp3", ".wav", ".flac", ".m4a",
    # 실행파일 (텍스트 추출 불가)
    ".exe", ".dll", ".so", ".bin",
    # 압축파일은 extract 단계에서 처리 — 실패 시 수동 검토를 위해 여기선 제외
    # .zip .tar .gz .7z .rar
    # 데이터 바이너리 (텍스트 추출 불가)
    ".db", ".sqlite", ".pkl", ".npy", ".npz",
    # 그래픽/컴파일 바이너리
    ".pbm", ".pgm", ".ppm", ".ai", ".o", ".obj", ".eps", ".ps",
}
# 텍스트 추출 가능한 문서 포맷 — 삭제하지 않고 AI가 처리
# .pdf .docx .docm .doc .xlsx .xls .pptx .ppt .hwp .hwpx
INDEX_FILE = WIKI_DIR / "index.md"
LOG_FILE = WIKI_DIR / "log.md"
SCHEMA_FILE = WIKI_DIR / "schema.md"
SKILL_FILE = Path(__file__).parents[1] / "SKILL.md"


def _parse_skill_help() -> str:
    """SKILL.md 서브명령어 테이블 파싱하여 help 출력"""
    if not SKILL_FILE.exists():
        return "SKILL.md를 찾을 수 없습니다."
    lines = SKILL_FILE.read_text(encoding="utf-8").splitlines()
    in_table = False
    rows = []
    for line in lines:
        if "| 명령어 | 설명 |" in line and "출력" in line:
            in_table = True
            rows.append(line)
            continue
        if in_table:
            if line.startswith("|"):
                rows.append(line)
            else:
                break
    return "\n".join(rows) if rows else "서브명령어 테이블을 찾을 수 없습니다."


def cmd_help():
    print("`oowiki help` 서브명령어 목록:\n")
    print(_parse_skill_help())


def cmd_version():
    print("oowiki 버전: v02")


def cmd_status():
    if not WIKI_DIR.exists():
        print(f"[oowiki status] 위키 폴더 없음: {WIKI_DIR}")
        print("초기화: 01_obsidian/0020_wiki/ 폴더를 생성하세요.")
        return

    md_files = [f for f in WIKI_DIR.rglob("*.md")
                if f.name not in ("index.md", "log.md", "schema.md")]
    categories = set(f.parent.name for f in md_files if f.parent != WIKI_DIR)

    # 마지막 수집일 (log.md에서)
    last_ingest = "없음"
    if LOG_FILE.exists():
        for line in LOG_FILE.read_text(encoding="utf-8").splitlines():
            if "ingest" in line.lower() and line.startswith("## "):
                last_ingest = line.replace("## ", "").split("—")[0].strip()

    print("[oowiki status]\n")
    print(f"위키 위치  : {WIKI_DIR}")
    print(f"총 페이지  : {len(md_files)}개")
    print(f"카테고리   : {len(categories)}개 ({', '.join(sorted(categories)) if categories else '없음'})")
    print(f"index.md   : {'[OK]' if INDEX_FILE.exists() else '[MISSING]'}")
    print(f"log.md     : {'[OK]' if LOG_FILE.exists() else '[MISSING]'}")
    print(f"schema.md  : {'[OK]' if SCHEMA_FILE.exists() else '[MISSING]'}")
    print(f"마지막 수집: {last_ingest}")


def cmd_list():
    if not WIKI_DIR.exists():
        print("위키 폴더가 없습니다.")
        return

    md_files = sorted(f for f in WIKI_DIR.rglob("*.md")
                      if f.name not in ("index.md", "log.md", "schema.md"))
    if not md_files:
        print("위키 페이지가 없습니다.")
        return

    print(f"[oowiki list] 총 {len(md_files)}개 페이지\n")
    current_cat = None
    for f in md_files:
        cat = f.parent.name if f.parent != WIKI_DIR else "(루트)"
        if cat != current_cat:
            print(f"\n### {cat}/")
            current_cat = cat
        print(f"  - {f.stem}")


def cmd_lint():
    if not WIKI_DIR.exists():
        print("위키 폴더가 없습니다.")
        return

    md_files = {f.stem: f for f in WIKI_DIR.rglob("*.md")
                if f.name not in ("index.md", "log.md", "schema.md")}

    issues = []

    # 1. index.md에 없는 고아 페이지
    if INDEX_FILE.exists():
        index_content = INDEX_FILE.read_text(encoding="utf-8")
        for stem, path in md_files.items():
            if f"[[{stem}]]" not in index_content and f"[[{stem}|" not in index_content:
                issues.append(f"[고아 페이지] {path.relative_to(WIKI_DIR)} — index.md에 미등록")
    else:
        issues.append("[구조 문제] index.md 없음")

    # 2. 끊어진 내부 링크
    for stem, path in md_files.items():
        content = path.read_text(encoding="utf-8")
        import re
        links = re.findall(r'\[\[([^\]|]+)', content)
        for link in links:
            link_stem = link.strip()
            if link_stem not in md_files:
                issues.append(f"[끊어진 링크] {stem} → [[{link_stem}]] 대상 없음")

    # 3. 빈 페이지
    for stem, path in md_files.items():
        content = path.read_text(encoding="utf-8").strip()
        if len(content) < 50:
            issues.append(f"[빈 페이지] {stem} — 내용이 너무 짧음 ({len(content)}자)")

    print(f"[oowiki lint] 점검 완료 — {len(md_files)}개 페이지\n")
    if issues:
        print(f"발견된 이슈: {len(issues)}개\n")
        for issue in issues:
            print(f"  {issue}")
    else:
        print("이슈 없음 ✅")


def cmd_index():
    """index.md 재생성"""
    if not WIKI_DIR.exists():
        print("위키 폴더가 없습니다.")
        return

    md_files = sorted(f for f in WIKI_DIR.rglob("*.md")
                      if f.name not in ("index.md", "log.md", "schema.md"))

    # 카테고리별 분류
    from collections import defaultdict
    cats = defaultdict(list)
    for f in md_files:
        cat = f.parent.name if f.parent != WIKI_DIR else "(루트)"
        cats[cat].append(f)

    lines = [
        "# Wiki Index",
        f"_마지막 갱신: {datetime.now().strftime('%Y-%m-%d %H:%M')}_",
        "",
    ]
    for cat in sorted(cats):
        lines.append(f"## {cat}")
        lines.append("| 페이지 | 요약 | 수정일 |")
        lines.append("|--------|------|--------|")
        for f in sorted(cats[cat]):
            mtime = datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d")
            # 첫 줄 요약 추출
            content = f.read_text(encoding="utf-8").splitlines()
            summary = ""
            in_frontmatter = False
            for line in content:
                line = line.strip()
                if line == "---":
                    in_frontmatter = not in_frontmatter
                    continue
                if in_frontmatter:
                    continue
                if line and not line.startswith("#") and not line.startswith("|"):
                    summary = line[:60]
                    break
            lines.append(f"| [[{f.stem}]] | {summary} | {mtime} |")
        lines.append("")

    INDEX_FILE.write_text("\n".join(lines), encoding="utf-8")
    print(f"[oowiki index] index.md 재생성 완료 — {len(md_files)}개 페이지 등록")


def cmd_inbox_list() -> list:
    """0019_정리 폴더의 처리 대기 파일 목록 출력"""
    if not INBOX_DIR.exists():
        print(f"[oowiki inbox] 정리 폴더 없음: {INBOX_DIR}")
        return []

    all_files = sorted(f for f in INBOX_DIR.rglob("*") if f.is_file())
    if not all_files:
        print("[oowiki inbox] 0019_정리 폴더가 비어 있습니다.")
        return []

    print(f"[oowiki inbox] 처리 대기 파일: {len(all_files)}개\n")
    for f in all_files:
        size = f.stat().st_size
        rel = f.relative_to(INBOX_DIR)
        print(f"  - {rel}  ({size:,}B)")
    return all_files


def cmd_inbox_extract():
    """압축 파일 해제 — 바이너리 purge 및 변환보다 먼저 실행"""
    if not INBOX_DIR.exists():
        return

    extracted, failed = [], []

    # ZIP
    for f in list(INBOX_DIR.rglob("*.zip")):
        try:
            import zipfile
            out_dir = f.parent / f.stem
            out_dir.mkdir(exist_ok=True)
            with zipfile.ZipFile(f) as zf:
                zf.extractall(out_dir)
            f.unlink()
            extracted.append(f.name)
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # TAR (tar / tar.gz / tgz / tar.bz2 / tar.xz)
    import tarfile as _tarfile
    for pattern in ["*.tar", "*.tar.gz", "*.tgz", "*.tar.bz2", "*.tar.xz"]:
        for f in list(INBOX_DIR.rglob(pattern)):
            try:
                stem = f.name.replace(".tar.gz", "").replace(".tar.bz2", "").replace(".tar.xz", "").replace(".tgz", "").replace(".tar", "")
                out_dir = f.parent / stem
                out_dir.mkdir(exist_ok=True)
                with _tarfile.open(f) as tf:
                    tf.extractall(out_dir)
                f.unlink()
                extracted.append(f.name)
            except Exception as e:
                failed.append(f"{f.name}: {e}")

    # GZ (단일 파일, tar.gz 제외)
    for f in list(INBOX_DIR.rglob("*.gz")):
        if not f.name.endswith(".tar.gz"):
            try:
                import gzip, shutil
                out = f.with_suffix("")
                with gzip.open(f) as gf, open(out, "wb") as of:
                    shutil.copyfileobj(gf, of)
                f.unlink()
                extracted.append(f.name)
            except Exception as e:
                failed.append(f"{f.name}: {e}")

    # 7Z (py7zr 선택적)
    for f in list(INBOX_DIR.rglob("*.7z")):
        try:
            import py7zr
            out_dir = f.parent / f.stem
            out_dir.mkdir(exist_ok=True)
            with py7zr.SevenZipFile(f) as sz:
                sz.extractall(out_dir)
            f.unlink()
            extracted.append(f.name)
        except ImportError:
            failed.append(f"{f.name}: py7zr 미설치 (uv add py7zr)")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # RAR (rarfile 선택적)
    for f in list(INBOX_DIR.rglob("*.rar")):
        try:
            import rarfile
            out_dir = f.parent / f.stem
            out_dir.mkdir(exist_ok=True)
            with rarfile.RarFile(f) as rf:
                rf.extractall(out_dir)
            f.unlink()
            extracted.append(f.name)
        except ImportError:
            failed.append(f"{f.name}: rarfile 미설치 (uv add rarfile)")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    total = len(extracted) + len(failed)
    if total == 0:
        return
    print(f"\n[oowiki extract] 압축 해제 {total}개")
    if extracted:
        names = ", ".join(extracted[:5]) + ("..." if len(extracted) > 5 else "")
        print(f"  완료 {len(extracted)}개: {names}")
    if failed:
        print(f"  실패/미지원 {len(failed)}개 (수동 검토 필요):")
        for msg in failed[:5]:
            print(f"    - {msg}")


def cmd_inbox_purge_binary() -> int:
    """0019_정리 폴더의 바이너리·이미지 파일을 자동 삭제하여 토큰 낭비 방지"""
    if not INBOX_DIR.exists():
        return 0

    deleted = []
    for f in INBOX_DIR.rglob("*"):
        if f.is_file() and f.suffix.lower() in BINARY_EXTENSIONS:
            f.unlink()
            deleted.append(f.name)

    if deleted:
        print(f"[oowiki inbox] 바이너리 파일 {len(deleted)}개 자동 삭제:")
        for name in deleted:
            print(f"  - {name}")
    return len(deleted)


def cmd_inbox_convert():
    """바이너리 문서를 텍스트로 변환 (wiki 통합 준비)"""
    if not INBOX_DIR.exists():
        return

    converted, failed, skipped = [], [], []

    # DOCX / DOCM → TXT
    for f in list(INBOX_DIR.rglob("*.docx")) + list(INBOX_DIR.rglob("*.docm")):
        try:
            from docx import Document
            doc = Document(f)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if text.strip():
                f.with_suffix(".txt").write_text(text, encoding="utf-8")
                f.unlink()
                converted.append(f.name)
            else:
                skipped.append(f"{f.name} (내용 없음)")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # XLSX → MD 표
    for f in INBOX_DIR.rglob("*.xlsx"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(f, read_only=True, data_only=True)
            lines = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                rows = [r for r in ws.iter_rows(values_only=True) if any(c is not None for c in r)]
                if not rows:
                    continue
                lines.append(f"## {sheet}")
                header = [str(c) if c is not None else "" for c in rows[0]]
                lines.append("| " + " | ".join(header) + " |")
                lines.append("|" + "|".join(["---"] * len(header)) + "|")
                for row in rows[1:]:
                    lines.append("| " + " | ".join(str(c) if c is not None else "" for c in row) + " |")
            if lines:
                f.with_suffix(".md").write_text("\n".join(lines), encoding="utf-8")
                f.unlink()
                converted.append(f.name)
            else:
                skipped.append(f"{f.name} (내용 없음)")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # PPTX → TXT
    for f in INBOX_DIR.rglob("*.pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(f)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
                if slide_texts:
                    lines.append(f"## Slide {i}")
                    lines.extend(slide_texts)
            if lines:
                f.with_suffix(".txt").write_text("\n".join(lines), encoding="utf-8")
                f.unlink()
                converted.append(f.name)
            else:
                skipped.append(f"{f.name} (내용 없음)")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # HWPX → TXT (ZIP + XML, 외부 라이브러리 불필요)
    for f in INBOX_DIR.rglob("*.hwpx"):
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            texts = []
            with zipfile.ZipFile(f) as zf:
                section_files = sorted(n for n in zf.namelist() if "Contents/section" in n and n.endswith(".xml"))
                for name in section_files:
                    with zf.open(name) as xf:
                        root = ET.parse(xf).getroot()
                        for elem in root.iter():
                            if elem.tag.endswith("}t") and elem.text and elem.text.strip():
                                texts.append(elem.text)
            if texts:
                f.with_suffix(".txt").write_text("\n".join(texts), encoding="utf-8")
                f.unlink()
                converted.append(f.name)
            else:
                skipped.append(f"{f.name} (내용 없음)")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # HWP → TXT (pyhwp CLI: hwp5txt)
    for f in INBOX_DIR.rglob("*.hwp"):
        try:
            import subprocess, shutil
            hwp5txt = shutil.which("hwp5txt")
            if not hwp5txt:
                failed.append(f"{f.name}: hwp5txt 미설치 (uv add pyhwp)")
                continue
            result = subprocess.run([hwp5txt, str(f)], capture_output=True, text=True, encoding="utf-8", errors="replace")
            text = result.stdout.strip()
            if text:
                f.with_suffix(".txt").write_text(text, encoding="utf-8")
                f.unlink()
                converted.append(f.name)
            else:
                failed.append(f"{f.name}: hwp5txt 출력 없음")
        except Exception as e:
            failed.append(f"{f.name}: {e}")

    # 지원 불가 포맷 안내
    unsupported = list(INBOX_DIR.rglob("*.doc")) + list(INBOX_DIR.rglob("*.xls")) + list(INBOX_DIR.rglob("*.ppt"))
    for f in unsupported:
        skipped.append(f"{f.name} (구버전 포맷 — LibreOffice 수동 변환 필요)")

    total = len(converted) + len(failed) + len(skipped)
    if total == 0:
        return
    print(f"\n[oowiki convert] 변환 대상 {total}개")
    if converted:
        print(f"  완료 {len(converted)}개: {', '.join(converted[:5])}{'...' if len(converted) > 5 else ''}")
    if skipped:
        print(f"  스킵 {len(skipped)}개")
    if failed:
        print(f"  실패 {len(failed)}개:")
        for msg in failed[:5]:
            print(f"    - {msg}")


def cmd_inbox_clear(targets: list | None = None):
    """0019_정리 폴더 내용물 삭제 (폴더 자체는 보존, targets=None이면 전체)"""
    if not INBOX_DIR.exists():
        print(f"[oowiki inbox] 정리 폴더 없음: {INBOX_DIR}")
        return

    if targets:
        files = [INBOX_DIR / t for t in targets]
    else:
        files = sorted(f for f in INBOX_DIR.rglob("*") if f.is_file())

    deleted = []
    failed = []
    for f in files:
        if f.exists() and f.is_file():
            try:
                f.unlink()
            except PermissionError:
                try:
                    import stat
                    f.chmod(stat.S_IWRITE)
                    f.unlink()
                except Exception:
                    failed.append(f.name)
                    continue
            deleted.append(f.name)

    # 빈 하위 폴더 정리 (INBOX_DIR 자체는 제외)
    for d in sorted(INBOX_DIR.rglob("*"), reverse=True):
        if d.is_dir() and d != INBOX_DIR and not any(d.iterdir()):
            try:
                d.rmdir()
            except Exception:
                pass

    print(f"[oowiki inbox] {len(deleted)}개 파일 삭제 완료 (0019_정리 폴더 유지)")
    if failed:
        print(f"  [경고] 삭제 실패 {len(failed)}개: {', '.join(failed[:5])}")
    for name in deleted:
        print(f"  - {name}")


def cmd_log_append(action: str, target: str, pages: list, summary: str):
    """log.md에 작업 기록 추가 (내부 유틸리티)"""
    WIKI_DIR.mkdir(parents=True, exist_ok=True)
    now = datetime.now().strftime("%Y-%m-%d %H:%M")
    entry = f"\n## {now} — {action}\n- **대상**: {target}\n- **결과**: {', '.join(pages) if pages else '없음'}\n- **요약**: {summary}\n"
    if LOG_FILE.exists():
        existing = LOG_FILE.read_text(encoding="utf-8")
        LOG_FILE.write_text(existing + entry, encoding="utf-8")
    else:
        LOG_FILE.write_text(f"# Wiki Log\n{entry}", encoding="utf-8")
    print(f"[log] {now} — {action} 기록 완료")


def main():
    args = sys.argv[1:]
    if not args or args[0] == "help":
        cmd_help()
    elif args[0] == "version":
        cmd_version()
    elif args[0] == "status":
        cmd_status()
    elif args[0] == "list":
        cmd_list()
    elif args[0] == "lint":
        cmd_lint()
    elif args[0] == "index":
        cmd_index()
    elif args[0] == "log":
        # log --append action target pages summary
        if len(args) >= 5:
            cmd_log_append(args[1], args[2], args[3].split(","), args[4])
        else:
            if LOG_FILE.exists():
                print(LOG_FILE.read_text(encoding="utf-8"))
            else:
                print("log.md 없음")
    elif args[0] == "run":
        print("[oowiki run] AI 주도 수집·통합을 시작합니다.")
        print("Claude가 index.md를 확인하고 위키에 통합합니다.")
        cmd_status()
        # 압축 파일 먼저 해제
        cmd_inbox_extract()
        # 바이너리 파일 자동 제거 (토큰 낭비 방지)
        cmd_inbox_purge_binary()
        # 문서 포맷 텍스트 변환 (docx/xlsx/pptx/hwp*)
        cmd_inbox_convert()
        # 0019_정리 inbox 파일 감지
        inbox_files = cmd_inbox_list()
        if inbox_files:
            print("\n[oowiki run] 위 inbox 파일들을 0020_wiki에 통합 후 'oowiki inbox --clear'로 삭제하세요.")
        # qmd 임베딩
        print("\n[oowiki run] qmd embed 실행 중...")
        import subprocess, shutil
        qmd_bin = shutil.which("qmd")
        if qmd_bin:
            cmd = [qmd_bin, "embed"]
        else:
            cmd = ["powershell", "-Command", "qmd embed"]
        result = subprocess.run(cmd, capture_output=False)
        if result.returncode != 0:
            print("[oowiki run] qmd embed 실패. 직접 실행: powershell -Command 'qmd embed'")
    elif args[0] == "extract":
        cmd_inbox_extract()
    elif args[0] == "convert":
        cmd_inbox_convert()
    elif args[0] == "inbox":
        if len(args) > 1 and args[1] == "--clear":
            targets = args[2:] if len(args) > 2 else None
            cmd_inbox_clear(targets)
        else:
            cmd_inbox_list()
    elif args[0] == "search":
        query = " ".join(args[1:]) if len(args) > 1 else ""
        print(f"[oowiki search] 질문: {query}")
        print("Claude가 index.md와 관련 페이지를 검색하여 답변합니다.")
        cmd_status()
    else:
        print(f"알 수 없는 명령: {args[0]}")
        cmd_help()


if __name__ == "__main__":
    main()
